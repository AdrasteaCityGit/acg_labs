import os
import zipfile
import tarfile
import rarfile
from datetime import datetime
import argparse
import json
import PyPDF2
from PIL import Image

# Import statements for Microsoft Office formats
try:
    import docx
    import openpyxl
    from pptx import Presentation
except ImportError:
    print("To extract metadata from Microsoft Office files, you need to install the required libraries.")
    print("You can install them using the following commands:")
    print("pip install python-docx openpyxl python-pptx")
    exit()

class ArchiveMetadataExtractor:
    SUPPORTED_EXTENSIONS = ['.zip', '.tar', '.rar', '.pdf', '.jpg', '.png', '.jpeg', '.docx', '.xlsx', '.pptx']

    def __init__(self, file_path):
        self.file_path = file_path
        self.metadata = {}

    def extract_metadata(self, selected_files=None):
        _, file_extension = os.path.splitext(self.file_path)

        if file_extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file format: {file_extension}")

        extraction_methods = {
            '.zip': self.extract_zip_metadata,
            '.tar': self.extract_tar_metadata,
            '.rar': self.extract_rar_metadata,
            '.pdf': self.extract_pdf_metadata,
            '.jpg': self.extract_image_metadata,
            '.png': self.extract_image_metadata,
            '.jpeg': self.extract_image_metadata,
            '.docx': self.extract_docx_metadata,
            '.xlsx': self.extract_xlsx_metadata,
            '.pptx': self.extract_pptx_metadata,
        }

        extraction_methods[file_extension](selected_files)

    def extract_zip_metadata(self, selected_files=None):
        with zipfile.ZipFile(self.file_path, 'r') as zip_file:
            self._extract_metadata_from_archive(zip_file, selected_files)

    def extract_tar_metadata(self, selected_files=None):
        with tarfile.open(self.file_path, 'r') as tar_file:
            self._extract_metadata_from_archive(tar_file, selected_files)

    def extract_rar_metadata(self, selected_files=None):
        with rarfile.RarFile(self.file_path, 'r') as rar_file:
            self._extract_metadata_from_archive(rar_file, selected_files)

    def _extract_metadata_from_archive(self, archive, selected_files=None):
        for info in archive.infolist():
            if not selected_files or info.filename in selected_files:
                self.metadata[info.filename] = self.get_common_metadata(info)
                if hasattr(info, 'compress_type'):
                    self.metadata[info.filename]['compression_type'] = info.compress_type

    def extract_pdf_metadata(self):
        with open(self.file_path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            info = pdf_reader.getDocumentInfo()
            self.metadata['PDF_Metadata'] = {
                'title': info.title,
                'author': info.author,
                'subject': info.subject,
                'producer': info.producer,
                'created_date': info.created.strftime('%Y-%m-%d %H:%M:%S UTC'),
                'modified_date': info.modified.strftime('%Y-%m-%d %H:%M:%S UTC'),
                'number_of_pages': len(pdf_reader.pages)
            }

    def extract_image_metadata(self):
        img = Image.open(self.file_path)
        self.metadata['Image_Metadata'] = {
            'format': img.format,
            'mode': img.mode,
            'size': img.size,
            'info': img.info
        }

    def extract_docx_metadata(self):
        doc = docx.Document(self.file_path)
        paragraphs = [paragraph.text for paragraph in doc.paragraphs]
        self.metadata['Docx_Metadata'] = {
            'paragraphs': paragraphs,
            'word_count': sum(len(paragraph.split()) for paragraph in paragraphs),
        }

    def extract_xlsx_metadata(self):
        wb = openpyxl.load_workbook(self.file_path)
        sheet_names = wb.sheetnames
        self.metadata['Xlsx_Metadata'] = {
            'sheet_names': sheet_names,
        }

    def extract_pptx_metadata(self):
        prs = Presentation(self.file_path)
        slide_titles = [slide.shapes.title.text for slide in prs.slides if slide.shapes.title]
        self.metadata['Pptx_Metadata'] = {
            'slide_titles': slide_titles,
        }

    @staticmethod
    def get_common_metadata(info):
        return {
            'file_size': info.size,
            'compression_type': None,
            'last_modified': datetime.utcfromtimestamp(info.mtime).strftime('%Y-%m-%d %H:%M:%S UTC')
        }

def display_metadata_summary(metadata):
    print("\nMetadata Summary:")
    print("----------------------------")
    for file_name, details in metadata.items():
        print(f"File: {file_name}")
        print(f"Size: {details['file_size']} bytes")
        print(f"Last Modified: {details['last_modified']}")
        if details.get('compression_type') is not None:
            print(f"Compression Type: {details['compression_type']}")
        print("----------------------------")

def save_metadata_to_json(metadata, output_file='metadata.json'):
    with open(output_file, 'w') as json_file:
        json.dump(metadata, json_file, indent=2)
    print(f"Metadata saved to {output_file}.")

def main():
    parser = argparse.ArgumentParser(description='Extract metadata from various file formats.')
    parser.add_argument('file_path', help='Path to the file or archive.')
    parser.add_argument('--json', action='store_true', help='Display metadata in JSON format.')
    parser.add_argument('--verbose', action='store_true', help='Display additional details.')
    parser.add_argument('--save-json', metavar='output_file', help='Save metadata to a JSON file.')
    parser.add_argument('--select-files', nargs='+', metavar='file_name',
                        help='Select specific file(s) within an archive to extract metadata.')

    args = parser.parse_args()
    file_path = args.file_path

    try:
        metadata_extractor = ArchiveMetadataExtractor(file_path)
        metadata_extractor.extract_metadata(args.select_files)

        if metadata_extractor.metadata:
            if args.verbose:
                print(f"Successfully extracted metadata from {file_path}.")
            display_metadata_summary(metadata_extractor.metadata)
            if args.json:
                print(json.dumps(metadata_extractor.metadata, indent=2))
            if args.save_json:
                save_metadata_to_json(metadata_extractor.metadata, args.save_json)

    except Exception as e:
        print(f"Error: {str(e)}")

if __name__ == "__main__":
    main()
